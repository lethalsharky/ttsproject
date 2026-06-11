import ollama
import asyncio
import edge_tts
import sounddevice as sd
import soundfile as sf
import threading
import queue as _q
import re
import os
import tempfile
import time

MODEL = "llama3.1:8b"
VOICE = "en-US-ChristopherNeural"

_gen_queue = _q.Queue()
_play_queue = _q.Queue()
_turn = [0]
_first_of_turn = set()
_bot_label_ready = threading.Event()
_bot_label_ready.set()  # pre-set so first input() shows immediately
_stop_text = threading.Event()

def clean(text):
    text = re.sub(r'\*+', '', text)
    text = re.sub(r'#+\s?', '', text)
    text = re.sub(r'`+', '', text)
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    text = re.sub(r'^\s*[-•]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

async def speak_async(text, path):
    communicate = edge_tts.Communicate(text, VOICE, rate="+0%", pitch="+2Hz")
    await communicate.save(path)

def _drain(q):
    while not q.empty():
        try: q.get_nowait()
        except _q.Empty: break

def _generator_worker():
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    while True:
        try:
            text, turn = _gen_queue.get(timeout=0.5)
        except _q.Empty:
            continue
        if text is None:
            _play_queue.put((None, None, turn))
            continue
        if turn != _turn[0]:
            continue
        tmp = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
        tmp.close()
        try:
            loop.run_until_complete(speak_async(text, tmp.name))
        except Exception:
            try: os.unlink(tmp.name)
            except OSError: pass
            continue
        if turn == _turn[0]:
            _play_queue.put((tmp.name, text, turn))
        else:
            try: os.unlink(tmp.name)
            except OSError: pass

def _player_worker():
    while True:
        try:
            path, text, turn = _play_queue.get(timeout=0.5)
        except _q.Empty:
            continue

        if path is None:
            _first_of_turn.discard(turn)
            _bot_label_ready.set()  # fallback if TTS failed for all sentences
            continue

        if turn != _turn[0]:
            try: os.unlink(path)
            except OSError: pass
            continue

        try:
            info = sf.info(path)
            duration = info.duration
            data, samplerate = sf.read(path)
        except Exception:
            try: os.unlink(path)
            except OSError: pass
            continue

        if turn not in _first_of_turn:
            _first_of_turn.add(turn)
            print('Bot: ', end='', flush=True)
            _bot_label_ready.set()

        char_delay = max(0.008, duration * 0.8 / max(len(text), 1))

        audio_done = threading.Event()

        def play_audio(data=data, samplerate=samplerate, turn=turn):
            sd.play(data, samplerate)
            try:
                while sd.get_stream().active:
                    if turn != _turn[0]:
                        sd.stop()
                        break
                    time.sleep(0.04)
            except Exception:
                pass
            audio_done.set()

        threading.Thread(target=play_audio, daemon=True).start()

        for ch in text + ' ':
            if turn != _turn[0] or _stop_text.is_set():
                break
            print(ch, end='', flush=True)
            time.sleep(char_delay)

        audio_done.wait(timeout=duration + 1.0)

        try: os.unlink(path)
        except OSError: pass

def speak_stop():
    _turn[0] += 1
    sd.stop()
    _drain(_gen_queue)
    _drain(_play_queue)
    _bot_label_ready.set()
    _stop_text.set()
    print('', flush=True)

def speak_sentence(text):
    t = clean(text)
    if t:
        _gen_queue.put((t, _turn[0]))

def read_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages).strip()

def read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"Slide {i}:"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
        if len(parts) > 1:
            slides.append("\n".join(parts))
    return "\n\n".join(slides).strip()

def load_file(path, history, ctx_holder):
    path = path.strip().strip('"').strip("'").strip()
    path = os.path.expandvars(os.path.expanduser(path))
    path = os.path.abspath(path)
    if not os.path.isfile(path):
        print(f"[File not found: {path}]\n")
        return
    try:
        if path.lower().endswith(".pdf"):
            content = read_pdf(path)
        elif path.lower().endswith(".pptx"):
            content = read_pptx(path)
        else:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
    except Exception as e:
        print(f"[Could not read file: {e}]\n")
        return
    name = os.path.basename(path)
    history[:] = [m for m in history if m.get("role") != "system"]
    history.insert(0, {
        "role": "system",
        "content": f'The user has loaded a file called "{name}". Its full contents are below. Answer all questions using this content.\n\n{content}'
    })
    chars = len(content)
    ctx_holder[0] = max(8192, (chars // 3) + 4096)
    print(f'[Loaded "{name}" — {chars} chars, ctx set to {ctx_holder[0]}]\n')

def chat():
    print("Voice Chatbot (type 'quit' to exit)\n")
    print("Commands: /load <path>  —  load a file into context\n")
    history = []
    ctx_holder = [4096]

    threading.Thread(target=_generator_worker, daemon=True).start()
    threading.Thread(target=_player_worker, daemon=True).start()

    while True:
        _bot_label_ready.wait()   # wait until "Bot: " is on screen
        time.sleep(0.35)          # let some chars appear first
        _stop_text.set()          # stop text printing, voice keeps playing
        time.sleep(0.025)         # let player see the flag and break
        print()                   # newline so "You: " is on its own line
        user_input = input("You: ").strip()
        speak_stop()

        if user_input.lower() in ("quit", "exit"):
            break
        if not user_input:
            continue

        if user_input.lower().startswith("/load "):
            load_file(user_input[6:], history, ctx_holder)
            continue

        # detect a file path anywhere in the message (e.g. "open C:\foo\bar.pdf summarize it")
        _skip_append = False
        _path_found = re.search(r'[A-Za-z]:\\[^\s"\'*?<>|]+|~?/[^\s"\'*?<>|]+', user_input)
        if _path_found:
            _candidate = os.path.abspath(os.path.expandvars(os.path.expanduser(_path_found.group(0))))
            if os.path.isfile(_candidate):
                load_file(_candidate, history, ctx_holder)
                _extra = user_input.replace(_path_found.group(0), '').strip().strip('"').strip("'").strip()
                if _extra and _extra.lower() not in ('open', 'load', 'read'):
                    history.append({"role": "user", "content": _extra})
                    _skip_append = True  # already appended, fall through to LLM
                else:
                    continue  # just loaded, no question yet

        if not _skip_append:
            history.append({"role": "user", "content": user_input})
        _bot_label_ready.clear()
        _stop_text.clear()

        full_reply = ""
        sentence_buf = ""

        for chunk in ollama.chat(
            model=MODEL,
            messages=history,
            stream=True,
            options={"num_gpu": 99, "num_ctx": ctx_holder[0], "temperature": 0.7},
        ):
            token = chunk["message"]["content"]
            full_reply += token
            sentence_buf += token

            while True:
                m = re.search(r'(?<=[.!?])\s', sentence_buf)
                if not m:
                    break
                speak_sentence(sentence_buf[:m.start()])
                sentence_buf = sentence_buf[m.end():]

        if sentence_buf.strip():
            speak_sentence(sentence_buf)

        _gen_queue.put((None, _turn[0]))
        history.append({"role": "assistant", "content": full_reply})

if __name__ == "__main__":
    chat()
